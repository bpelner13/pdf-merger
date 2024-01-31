#!/usr/bin/ python3
"""Got sick of merging powerpoints too..."""
from tkinter import Tk
from tkinter.filedialog import askopenfilenames

import pptx


def getFiles():
    """Get names of files user wants to merge"""
    Tk().withdraw()
    files = askopenfilenames()
    return files


def mergePpts(ppts):
    """Merge powerpoints together"""
    filename = "merged.pptx"
    pp1 = pptx.Presentation(ppts[0])
    for pp in ppts[1:]:
        p = pptx.Presentation(pp)
        for slide in p.slides:
            sl = pp1.slides.add_slide(pp1.slide_layouts[1])
            sl.shapes.title.text = slide.shapes.title.text
            try:
                sl.placeholders[1].text = slide.placeholders[1].text
            except:
                sl.placeholders[1].image = slide.placeholders[1].image
    pp1.save(filename)
    print("Created: :", filename)



if __name__ == '__main__':
    files = getFiles()
    print("Merging: ", *files)
    mergePpts(files)
